import json
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PresentationMaker:
    def __init__(self):
        self._templates = {}

    def add_template(self, name, slides):
        self._templates[name] = slides

    def get_template(self, name):
        template = self._templates.get(name)
        if not template:
            raise KeyError(f"Template '{name}' not found")
        return template

    def list_templates(self):
        return list(self._templates.keys())

    def remove_template(self, name):
        if name not in self._templates:
            raise KeyError(f"Template '{name}' not found")
        del self._templates[name]

    def create_from_outline(self, outline, filepath):
        if not HAS_PPTX:
            raise NotImplementedError("python-pptx is required. Install with: pip install python-pptx")
        prs = Presentation()
        if isinstance(outline, str):
            outline = self._parse_outline_text(outline)
        for slide_data in outline:
            slide_layout = prs.slide_layouts.get(slide_data.get("layout", 1))
            if slide_layout is None:
                slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            if title_shape and "title" in slide_data:
                title_shape.text = slide_data["title"]
            if "content" in slide_data:
                content = slide_data["content"]
                if slide.placeholders:
                    body_shape = slide.placeholders[1]
                    body_shape.text = content if isinstance(content, str) else "\n".join(content)
            if "bullets" in slide_data:
                if slide.placeholders:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for i, bullet in enumerate(slide_data["bullets"]):
                        if i == 0:
                            tf.text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
        prs.save(filepath)
        return filepath

    def create_from_template(self, template_name, filepath, **kwargs):
        template = self.get_template(template_name)
        filled = []
        for slide_data in template:
            filled_slide = dict(slide_data)
            for key, value in slide_data.items():
                if isinstance(value, str):
                    filled_slide[key] = value.format(**kwargs)
            filled.append(filled_slide)
        return self.create_from_outline(filled, filepath)

    def _parse_outline_text(self, text):
        slides = []
        current_slide = None
        for line in text.strip().split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("## "):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": line[3:], "bullets": []}
            elif line.startswith("- ") and current_slide:
                current_slide["bullets"].append(line[2:])
            elif current_slide:
                current_slide["bullets"].append(line)
        if current_slide:
            slides.append(current_slide)
        return slides

    def add_slide_numbering(self, filepath):
        if not HAS_PPTX:
            raise NotImplementedError("python-pptx is required")
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame and "{{slide_number}}" in shape.text:
                    shape.text = shape.text.replace("{{slide_number}}", str(i))
        prs.save(filepath)

    def merge_presentations(self, filepaths, output_path):
        if not HAS_PPTX:
            raise NotImplementedError("python-pptx is required")
        prs = Presentation()
        for filepath in filepaths:
            source = Presentation(filepath)
            for slide in source.slides:
                slide_layout = prs.slide_layouts.get(1)
                if slide_layout is None:
                    slide_layout = prs.slide_layouts[0]
                new_slide = prs.slides.add_slide(slide_layout)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        new_shape = new_slide.shapes.add_textbox(
                            shape.left, shape.top, shape.width, shape.height
                        )
                        new_shape.text_frame.text = shape.text_frame.text
        prs.save(output_path)
        return output_path
